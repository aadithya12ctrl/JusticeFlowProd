"""
JusticeFlow Pitch Deck Generator — Autumn Glassmorphic Theme
Generates a 12-slide PPTX using python-pptx 1.0.2
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ═══════════════════════════════════════════════
# THEME
# ═══════════════════════════════════════════════
BG_DARK       = RGBColor(0x2C, 0x1A, 0x0E)
AMBER         = RGBColor(0xD4, 0x78, 0x2A)
SOFT_AMBER    = RGBColor(0xF4, 0xA2, 0x61)
WHEAT         = RGBColor(0xF5, 0xDE, 0xB3)
ESPRESSO      = RGBColor(0x3E, 0x1C, 0x00)
WOOD          = RGBColor(0x8B, 0x5A, 0x2B)
TAUPE         = RGBColor(0x9E, 0x6B, 0x47)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
SIENNA        = RGBColor(0xA0, 0x52, 0x2D)
SAGE          = RGBColor(0x6B, 0x8F, 0x71)
DARK_RED      = RGBColor(0xB0, 0x3A, 0x2E)
DARK_WOOD     = RGBColor(0x5C, 0x40, 0x33)


def _set_alpha(shape, alpha_pct):
    """Set fill transparency via direct XML manipulation. alpha_pct=15 means 15% opaque (85% transparent)."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        return
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    # Remove existing alpha
    for old in srgb.findall(qn('a:alpha')):
        srgb.remove(old)
    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
    alpha_el.set('val', str(int(alpha_pct * 1000)))  # 15000 = 15%


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def glass_card(slide, x, y, w, h, border=AMBER, opacity_pct=18):
    """Glassmorphic card — white fill with low opacity & amber border."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    _set_alpha(s, opacity_pct)
    s.line.color.rgb = border
    s.line.width = Pt(1)
    return s


def accent_oval(slide, x, y, w, h, color, opacity_pct=30):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    _set_alpha(s, opacity_pct)
    s.line.fill.background()
    return s


def leaf_accents(slide):
    accent_oval(slide, Inches(9.2), Inches(0.15), Inches(0.4), Inches(0.4), AMBER, 30)
    accent_oval(slide, Inches(0.2), Inches(4.95), Inches(0.35), Inches(0.35), SIENNA, 25)


def bottom_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.3), Inches(10), Inches(0.325))
    s.fill.solid()
    s.fill.fore_color.rgb = AMBER
    s.line.fill.background()


def txt(slide, text, x, y, w, h, name="Georgia", sz=40, bold=True, italic=False,
        color=WHEAT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = name
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tb


def bullet_card(slide, x, y, w, h, title, items, title_clr=AMBER,
                text_clr=ESPRESSO, border=AMBER):
    """Glass card with a title + bullet items."""
    glass_card(slide, x, y, w, h, border)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = title_clr
    p.space_after = Pt(6)
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Calibri"
        p.font.size = Pt(10.5)
        p.font.color.rgb = text_clr
        p.space_after = Pt(3)


# ═══════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════

def s01_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    accent_oval(sl, Inches(3.5), Inches(0.3), Inches(3), Inches(3), AMBER, 18)
    txt(sl, "⚖️ JusticeFlow", Inches(1), Inches(1.3), Inches(8), Inches(1.2),
        "Georgia", 52, True, True, WHEAT, PP_ALIGN.CENTER)
    txt(sl, "AI-Powered Dispute Resolution Platform", Inches(1), Inches(2.3),
        Inches(8), Inches(0.6), "Calibri", 20, False, False, SOFT_AMBER, PP_ALIGN.CENTER)
    txt(sl, "Reducing India's 4.7 Crore case backlog through intelligent legal AI",
        Inches(1.5), Inches(2.9), Inches(7), Inches(0.5),
        "Calibri", 14, False, True, TAUPE, PP_ALIGN.CENTER)
    bottom_bar(sl)
    badges = ["Python 3.12", "Streamlit", "LangChain", "Groq LLaMA 3.3"]
    for i, b in enumerate(badges):
        bx = Inches(1.5 + i * 1.85)
        glass_card(sl, bx, Inches(3.7), Inches(1.6), Inches(0.5), AMBER, 22)
        txt(sl, b, bx + Inches(0.05), Inches(3.78), Inches(1.5), Inches(0.35),
            "Calibri", 11, True, False, ESPRESSO, PP_ALIGN.CENTER)
    leaf_accents(sl)


def s02_problem(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "The Crisis", Inches(0.5), Inches(0.2), Inches(5), Inches(0.65),
        "Georgia", 40, True, True, WHEAT)
    stats = [
        "4.7 Crore+ cases pending across Indian courts",
        "3–5 years average civil dispute resolution time",
        "21 judges per million citizens (lowest in democracies)",
        "70% of cases are settleable before trial (NALSA)",
        "Judges spend 40%+ time reviewing filterable cases",
        "Zero institutional memory across courthouses",
    ]
    bullet_card(sl, Inches(0.4), Inches(1.0), Inches(5.5), Inches(3.7),
                "🇮🇳 India's Judicial Backlog", stats)
    glass_card(sl, Inches(6.2), Inches(1.0), Inches(3.3), Inches(3.7))
    txt(sl, "4.7Cr+", Inches(6.3), Inches(1.3), Inches(3.1), Inches(0.9),
        "Georgia", 48, True, False, SOFT_AMBER, PP_ALIGN.CENTER)
    txt(sl, "Pending Cases", Inches(6.3), Inches(2.1), Inches(3.1), Inches(0.4),
        "Calibri", 12, False, False, WOOD, PP_ALIGN.CENTER)
    txt(sl, "21", Inches(6.3), Inches(2.7), Inches(3.1), Inches(0.8),
        "Georgia", 44, True, False, SOFT_AMBER, PP_ALIGN.CENTER)
    txt(sl, "Judges per Million", Inches(6.3), Inches(3.4), Inches(3.1), Inches(0.4),
        "Calibri", 12, False, False, WOOD, PP_ALIGN.CENTER)
    txt(sl, "Source: National Judicial Data Grid (NJDG), 2024",
        Inches(0.4), Inches(5.0), Inches(5), Inches(0.3), "Calibri", 9, False, True, TAUPE)
    leaf_accents(sl)
    bottom_bar(sl)


def s03_solution(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Our Solution", Inches(0.5), Inches(0.1), Inches(5), Inches(0.6),
        "Georgia", 40, True, True, WHEAT)
    feats = [
        ("🧬", "Case DNA", "6D fingerprint + precedent twin matching"),
        ("🔍", "Auto Filter", "5-step hybrid rule + LLM screening"),
        ("📊", "DLS Engine", "Dismissal probability with per-reason breakdown"),
        ("🤝", "AI Negotiation", "Dual adversarial agents + live judge control"),
        ("📈", "Risk Indicator", "Linguistic risk scoring with cooling-off alerts"),
        ("🕸️", "Conflict Graph", "Entity network detects repeat offenders"),
        ("⚖️", "Judge Cockpit", "90-second review with downloadable PDF brief"),
    ]
    pos = [(0.3, 0.85), (3.5, 0.85), (6.7, 0.85),
           (0.3, 2.15), (3.5, 2.15), (6.7, 2.15),
           (3.5, 3.45)]
    for i, (icon, title, desc) in enumerate(feats):
        px, py = pos[i]
        glass_card(sl, Inches(px), Inches(py), Inches(3.0), Inches(1.15))
        txt(sl, f"{icon}  {title}", Inches(px + 0.12), Inches(py + 0.08),
            Inches(2.7), Inches(0.4), "Calibri", 14, True, False, AMBER)
        txt(sl, desc, Inches(px + 0.12), Inches(py + 0.5),
            Inches(2.7), Inches(0.5), "Calibri", 10.5, False, False, ESPRESSO)
    leaf_accents(sl)
    bottom_bar(sl)


def s03b_judging(prs):
    """SLIDE — Judging Criteria Alignment (10 criteria, 2 columns of 5)"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "How We Score", Inches(0.5), Inches(0.05), Inches(5), Inches(0.5),
        "Georgia", 36, True, True, WHEAT)
    txt(sl, "100 Marks  •  10 Judging Criteria  •  10 Each", Inches(4.5), Inches(0.12),
        Inches(5), Inches(0.35), "Calibri", 11, False, True, TAUPE, PP_ALIGN.RIGHT)

    criteria = [
        ("Social Impact",       "4.7Cr backlog → faster justice for underserved"),
        ("Legal Novelty",       "Multi-agent negotiation + live judge intervention"),
        ("Technical Feasibility", "Working prototype, ₹0 infra, demo-ready now"),
        ("Practicability",      "Indian statutes, ₹ amounts, Hindi toggle, tribunals"),
        ("Business Model",      "₹3K/seat B2G SaaS, 73% margin, Mo 10 break-even"),
        ("Scalability",         "5 districts → 50 → 200, data flywheel moat"),
        ("Target Market",       "District judges (30-40 cases/day) + legal aid lawyers"),
        ("Problem Solving",     "45 min → 90 sec review. 3-5 yrs → days settlement"),
        ("Creativity",          "Judge types mid-negotiation, AI agents react live"),
        ("Q&A Readiness",       "Precedent-backed claims, CPC Rule 11, RERA §18"),
    ]
    for i, (criterion, evidence) in enumerate(criteria):
        col = i % 2
        row = i // 2
        px = 0.3 + col * 4.9
        py = 0.6 + row * 0.92
        glass_card(sl, Inches(px), Inches(py), Inches(4.65), Inches(0.8))
        txt(sl, criterion, Inches(px + 0.12), Inches(py + 0.05),
            Inches(2.0), Inches(0.35), "Calibri", 12, True, False, AMBER)
        txt(sl, evidence, Inches(px + 0.12), Inches(py + 0.38),
            Inches(4.3), Inches(0.35), "Calibri", 9.5, False, False, ESPRESSO)
    leaf_accents(sl)
    bottom_bar(sl)


def s04_dna(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Case DNA Fingerprinting", Inches(0.5), Inches(0.2), Inches(8), Inches(0.65),
        "Georgia", 38, True, True, WHEAT)
    dims = [
        "📂  Category — RERA, Consumer, Labour, Motor Accident…",
        "🏛️  Jurisdiction — Court/tribunal clarity score",
        "💰  Claim Bucket — <₹50K | ₹50K–5L | ₹5L–50L | >₹50L",
        "📋  Evidence Strength — Documentation quality (0–1)",
        "💭  Emotional Intensity — Language heat score (0–1)",
        "🌟  Novelty — How uncommon the legal question is (0–1)",
    ]
    bullet_card(sl, Inches(0.3), Inches(1.0), Inches(5.0), Inches(3.2), "6 DNA Dimensions", dims)
    # Visual — overlapping ovals as radar proxy
    colors = [AMBER, SOFT_AMBER, SIENNA, WOOD, SAGE, DARK_WOOD]
    centers = [(7.5, 1.6), (8.4, 2.1), (8.1, 3.0), (7.2, 3.0), (6.8, 2.1), (7.5, 2.4)]
    labels = ["Cat", "Jur", "Claim", "Evid", "Emot", "Nov"]
    for (cx, cy), col, lbl in zip(centers, colors, labels):
        accent_oval(sl, Inches(cx), Inches(cy), Inches(0.85), Inches(0.85), col, 40)
        txt(sl, lbl, Inches(cx + 0.12), Inches(cy + 0.22), Inches(0.6), Inches(0.3),
            "Calibri", 9, True, False, WHEAT, PP_ALIGN.CENTER)
    glass_card(sl, Inches(0.3), Inches(4.4), Inches(9.4), Inches(0.55), AMBER, 20)
    txt(sl, "Cosine similarity matching against 12 real SC/HC precedents  •  Twin identification with % match",
        Inches(0.5), Inches(4.48), Inches(9), Inches(0.35), "Calibri", 11, False, True, ESPRESSO, PP_ALIGN.CENTER)
    leaf_accents(sl)
    bottom_bar(sl)


def s05_dls(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Dismissal Likelihood Score", Inches(0.5), Inches(0.2), Inches(8), Inches(0.65),
        "Georgia", 38, True, True, WHEAT)
    risks = [
        ("🏛️", "Jurisdiction", "Wrong court\nor forum"),
        ("⏰", "Limitation", "Filing deadline\nexpired"),
        ("📋", "Evidence", "Documentation\ngap found"),
        ("❌", "Frivolous", "Lacks legal\nmerit"),
        ("⚠️", "Procedural", "Process or\nformat error"),
    ]
    for i, (icon, title, desc) in enumerate(risks):
        rx = Inches(0.3 + i * 1.9)
        glass_card(sl, rx, Inches(1.1), Inches(1.7), Inches(2.1))
        txt(sl, icon, rx + Inches(0.3), Inches(1.2), Inches(1.1), Inches(0.5),
            "Calibri", 28, False, False, AMBER, PP_ALIGN.CENTER)
        txt(sl, title, rx + Inches(0.1), Inches(1.7), Inches(1.5), Inches(0.35),
            "Calibri", 13, True, False, AMBER, PP_ALIGN.CENTER)
        txt(sl, desc, rx + Inches(0.1), Inches(2.15), Inches(1.5), Inches(0.7),
            "Calibri", 10, False, False, ESPRESSO, PP_ALIGN.CENTER)
    glass_card(sl, Inches(0.5), Inches(3.5), Inches(9.0), Inches(0.85), DARK_RED, 20)
    txt(sl, "⚠️ Cases >75% DLS → mandatory judge override with full audit trail",
        Inches(0.7), Inches(3.58), Inches(8.5), Inches(0.35),
        "Calibri", 13, True, False, DARK_RED, PP_ALIGN.CENTER)
    txt(sl, "Based on CPC Order VII Rule 11 grounds for dismissal under Indian law",
        Inches(0.7), Inches(3.95), Inches(8.5), Inches(0.3),
        "Calibri", 10, False, True, WOOD, PP_ALIGN.CENTER)
    leaf_accents(sl)
    bottom_bar(sl)


def s06_negotiation(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "AI Negotiation Sandbox", Inches(0.5), Inches(0.2), Inches(8), Inches(0.65),
        "Georgia", 38, True, True, WHEAT)
    p_items = [
        "Argues maximum recovery for plaintiff",
        "Cites RERA §18, CPA 2019, ICA §73",
        "Starts strong, then strategically concedes",
        "Tracks fair settlement threshold (15%)",
    ]
    bullet_card(sl, Inches(0.3), Inches(1.0), Inches(4.2), Inches(2.3),
                "⚖️ Plaintiff Advocate AI", p_items, SAGE)
    # Divider
    div = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.85), Inches(1.0), Inches(0.05), Inches(2.3))
    div.fill.solid(); div.fill.fore_color.rgb = AMBER; div.line.fill.background()
    txt(sl, "VS", Inches(4.55), Inches(2.0), Inches(0.65), Inches(0.4),
        "Georgia", 16, True, False, SOFT_AMBER, PP_ALIGN.CENTER)
    d_items = [
        "Minimizes liability while staying reasonable",
        "Challenges evidence strength & jurisdiction",
        "Proposes counter-settlements with precedent",
        "Accepts when offer reaches fair range",
    ]
    bullet_card(sl, Inches(5.3), Inches(1.0), Inches(4.2), Inches(2.3),
                "🛡️ Defendant Advocate AI", d_items, SIENNA)
    glass_card(sl, Inches(0.5), Inches(3.55), Inches(9.0), Inches(1.15), AMBER)
    txt(sl, "👨‍⚖️  Human Judge Presides — Live Intervention",
        Inches(0.7), Inches(3.6), Inches(8.5), Inches(0.4),
        "Calibri", 15, True, False, AMBER, PP_ALIGN.CENTER)
    txt(sl, "Judge types directions mid-round  •  Force settlement or escalate  •  Downloadable settlement agreement",
        Inches(0.7), Inches(4.05), Inches(8.5), Inches(0.4),
        "Calibri", 11, False, False, ESPRESSO, PP_ALIGN.CENTER)
    leaf_accents(sl)
    bottom_bar(sl)


def s07_risk(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Litigation Risk Indicator", Inches(0.5), Inches(0.2), Inches(8), Inches(0.65),
        "Georgia", 38, True, True, WHEAT)
    levels = [
        ("0–20", "Calm", "Factual, measured language", SAGE),
        ("20–40", "Mild", "Some emotional undertones", RGBColor(0xE8, 0xD4, 0xA0)),
        ("40–60", "Elevated", "Frustration evident in filing", SOFT_AMBER),
        ("60–80", "High", "Hostile language detected", AMBER),
        ("80–100", "Critical", "Threats or extreme distress", DARK_RED),
    ]
    for i, (score, level, desc, bg) in enumerate(levels):
        ry = Inches(1.05 + i * 0.6)
        bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), ry, Inches(6.5), Inches(0.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = bg
        _set_alpha(bar, 40)
        bar.line.color.rgb = bg; bar.line.width = Pt(0.5)
        tc = WHITE if bg in (DARK_RED, AMBER) else ESPRESSO
        txt(sl, score, Inches(0.65), ry + Inches(0.05), Inches(0.8), Inches(0.35), "Calibri", 11, True, False, tc)
        txt(sl, level, Inches(1.6), ry + Inches(0.05), Inches(1.2), Inches(0.35), "Calibri", 11, True, False, tc)
        txt(sl, desc, Inches(3.0), ry + Inches(0.05), Inches(3.8), Inches(0.35), "Calibri", 10.5, False, False, tc)
    glass_card(sl, Inches(7.3), Inches(1.05), Inches(2.3), Inches(3.1))
    txt(sl, "🔥", Inches(7.5), Inches(1.15), Inches(1.9), Inches(0.5), "Calibri", 36, False, False, AMBER, PP_ALIGN.CENTER)
    txt(sl, "Cooling-off trigger\nat score 70+", Inches(7.4), Inches(1.7), Inches(2.1), Inches(0.7),
        "Calibri", 12, True, False, AMBER, PP_ALIGN.CENTER)
    txt(sl, "3× failure rate when emotional intensity exceeds threshold.\nAuto-routes to Lok Adalat.",
        Inches(7.4), Inches(2.5), Inches(2.1), Inches(1.0), "Calibri", 9, False, False, ESPRESSO, PP_ALIGN.CENTER)
    leaf_accents(sl)
    bottom_bar(sl)


def s08_audience(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Who Is This For?", Inches(0.5), Inches(0.15), Inches(5), Inches(0.6),
        "Georgia", 38, True, True, WHEAT)
    audiences = [
        ("👨‍⚖️ Judges & Magistrates", ["90-second review via Judge Cockpit", "AI-generated statute citations", "One-click downloadable PDF brief"]),
        ("👔 Legal Representatives", ["Case DNA matching for precedent research", "Negotiation prep with AI sandbox", "DLS risk assessment for clients"]),
        ("🏛️ Court Administrators", ["Conflict graph for case distribution", "Auto-filter reduces trivial filings", "Analytics dashboard for court KPIs"]),
        ("👥 Litigants & Public", ["Hindi language support (हिंदी)", "Faster pre-trial settlement via AI", "Transparent, auditable AI process"]),
    ]
    pos = [(0.3, 0.85), (5.1, 0.85), (0.3, 2.95), (5.1, 2.95)]
    for (title, items), (px, py) in zip(audiences, pos):
        bullet_card(sl, Inches(px), Inches(py), Inches(4.5), Inches(1.85), title, items)
    leaf_accents(sl)
    bottom_bar(sl)


def s09_tech(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Architecture & Tech Stack", Inches(0.5), Inches(0.15), Inches(8), Inches(0.6),
        "Georgia", 38, True, True, WHEAT)
    stack = [
        ("Groq Cloud — LLaMA 3.3 70B", AMBER),
        ("LangChain LCEL Pipelines", SOFT_AMBER),
        ("Streamlit Frontend", SIENNA),
        ("Python 3.12 Runtime", WOOD),
        ("SQLite + NetworkX", DARK_WOOD),
        ("Plotly Visualizations", TAUPE),
    ]
    for i, (label, color) in enumerate(stack):
        sy = Inches(0.9 + i * 0.65)
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), sy, Inches(4.0), Inches(0.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        txt(sl, label, Inches(0.55), sy + Inches(0.05), Inches(3.7), Inches(0.35),
            "Calibri", 12, True, False, WHEAT)
    pipeline = [
        "Case Filing → DNA Vector Computed",
        "DLS Analysis → Risk Scores Generated",
        "Auto-Filter → 5 Rule + LLM Checks",
        "AI Negotiation → Settlement Attempted",
        "Judge Cockpit → 90-Second Review",
        "Conflict Graph → Pattern Detection",
    ]
    bullet_card(sl, Inches(4.8), Inches(0.9), Inches(4.8), Inches(3.0),
                "🔄 AI Pipeline Flow", pipeline)
    glass_card(sl, Inches(4.8), Inches(4.1), Inches(4.8), Inches(0.6), SAGE, 22)
    txt(sl, "💰 Infrastructure Cost: ₹0 — Free-tier APIs, no GPU required",
        Inches(5.0), Inches(4.18), Inches(4.5), Inches(0.35),
        "Calibri", 11.5, True, False, SAGE, PP_ALIGN.CENTER)
    leaf_accents(sl)
    bottom_bar(sl)


def s10_business(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Business Model", Inches(0.5), Inches(0.15), Inches(8), Inches(0.6),
        "Georgia", 38, True, True, WHEAT)
    tiers = [
        ("🏛️ Court License (B2G)", ["₹3,000/judge/month", "₹5,000/court admin/month", "₹2,00,000 one-time setup"], AMBER),
        ("👔 Lawyer SaaS (B2B)", ["₹999/mo Solo practitioner", "₹7,999/mo Law firm (10 seats)", "Custom enterprise pricing"], SOFT_AMBER),
        ("📊 Data & Insights API", ["₹50,000/month analytics API", "Policy reports & trend data", "Anonymized court analytics"], SIENNA),
    ]
    for i, (title, items, accent) in enumerate(tiers):
        bullet_card(sl, Inches(0.3), Inches(0.85 + i * 1.45), Inches(6.5), Inches(1.3),
                    title, items, accent)
    glass_card(sl, Inches(7.1), Inches(0.85), Inches(2.5), Inches(3.8))
    metrics = [("73%", "Gross Margin"), ("68:1", "LTV:CAC"), ("Mo 10", "Break-Even"), ("5×", "Yr 2 Return")]
    for i, (val, label) in enumerate(metrics):
        my = Inches(1.0 + i * 0.85)
        txt(sl, val, Inches(7.2), my, Inches(2.3), Inches(0.5),
            "Georgia", 24, True, False, SOFT_AMBER, PP_ALIGN.CENTER)
        txt(sl, label, Inches(7.2), my + Inches(0.4), Inches(2.3), Inches(0.3),
            "Calibri", 10, False, False, WOOD, PP_ALIGN.CENTER)
    leaf_accents(sl)
    bottom_bar(sl)


def s11_financials(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "Growth Projections", Inches(0.5), Inches(0.15), Inches(8), Inches(0.6),
        "Georgia", 38, True, True, WHEAT)
    from pptx.chart.data import CategoryChartData
    cd = CategoryChartData()
    cd.categories = ['Year 1', 'Year 2', 'Year 3']
    cd.add_series('Revenue (₹ Cr)', (1.1, 8.5, 25.0))
    cf = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                             Inches(0.5), Inches(0.9), Inches(5.5), Inches(3.5), cd)
    chart = cf.chart
    chart.has_legend = False
    series = chart.plots[0].series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = AMBER
    series.data_labels.show_value = True
    series.data_labels.font.size = Pt(14)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = WHEAT
    series.data_labels.number_format = u'₹#.0"Cr"'
    chart.plots[0].gap_width = 120
    milestones = [
        ("Year 1", "5 districts • 500 users\nSeed Stage"),
        ("Year 2", "50 districts • 5K users\nSeries A"),
        ("Year 3", "200 districts • 20K users\nSeries B"),
    ]
    for i, (title, desc) in enumerate(milestones):
        my = Inches(0.9 + i * 1.4)
        glass_card(sl, Inches(6.3), my, Inches(3.3), Inches(1.15))
        txt(sl, title, Inches(6.45), my + Inches(0.05), Inches(3.0), Inches(0.3),
            "Calibri", 14, True, False, AMBER)
        txt(sl, desc, Inches(6.45), my + Inches(0.4), Inches(3.0), Inches(0.65),
            "Calibri", 10, False, False, ESPRESSO)
    leaf_accents(sl)
    bottom_bar(sl)


def s12_ask(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    txt(sl, "The Ask & Roadmap", Inches(0.5), Inches(0.15), Inches(8), Inches(0.6),
        "Georgia", 38, True, True, WHEAT)
    invest = [
        "Karnataka pilot deployment — ₹15L",
        "Engineering team (2 devs × 6mo) — ₹18L",
        "Govt relations & compliance — ₹7L",
        "Cloud infrastructure & APIs — ₹5L",
        "Marketing & legal conferences — ₹5L",
    ]
    bullet_card(sl, Inches(0.3), Inches(0.85), Inches(4.6), Inches(2.8),
                "💰 Seeking ₹50 Lakh Seed Investment", invest)
    txt(sl, "₹50L for 10% equity  •  Break-even Month 10  •  5× return by Year 2",
        Inches(0.4), Inches(3.7), Inches(4.4), Inches(0.35),
        "Calibri", 10, True, False, SOFT_AMBER, PP_ALIGN.CENTER)
    roadmap = [
        "RAG with Indian Kanoon / SCC Online",
        "Bhashini API — Hindi, Kannada, Tamil",
        "Voice input via Whisper API",
        "Appeal Risk Predictor module",
        "Fine-tuned Indian Legal LLM",
        "Mobile App (Flutter)",
    ]
    bullet_card(sl, Inches(5.2), Inches(0.85), Inches(4.4), Inches(3.2),
                "🚀 Future Roadmap", roadmap)
    txt(sl, "\"4.7 crore cases are waiting. We can give every judge an AI co-pilot.\"",
        Inches(0.5), Inches(4.35), Inches(9.0), Inches(0.5),
        "Georgia", 16, True, True, SOFT_AMBER, PP_ALIGN.CENTER)
    leaf_accents(sl)
    bottom_bar(sl)


# ═══════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    s01_title(prs)
    s02_problem(prs)
    s03_solution(prs)
    s03b_judging(prs)
    s04_dna(prs)
    s05_dls(prs)
    s06_negotiation(prs)
    s07_risk(prs)
    s08_audience(prs)
    s09_tech(prs)
    s10_business(prs)
    s11_financials(prs)
    s12_ask(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "JusticeFlow_Pitch_Deck_v2.pptx")
    prs.save(out)
    print(f"✅ Saved: {out}")
    print(f"   {len(prs.slides)} slides generated")

if __name__ == "__main__":
    main()
